"""Build PowerPoint presentation from v6 report content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRY = RGBColor(0xCC, 0xD6, 0xE0)
GREEN     = RGBColor(0x06, 0xD6, 0x8A)
ORANGE    = RGBColor(0xFF, 0x9F, 0x1C)
SLIDE_W   = Inches(13.33)
SLIDE_H   = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_paragraph(tf, text, font_size=14, bold=False,
                  color=WHITE, align=PP_ALIGN.LEFT,
                  space_before=Pt(4), bullet_char="▸ "):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = (bullet_char if bullet_char else "") + text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def bg(slide):
    """Fill slide background with dark colour."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)


def header_bar(slide, title_text, sub_text=""):
    """Top accent bar with section title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), RGBColor(0x05, 0x28, 0x45))
    # cyan left strip
    add_rect(slide, 0, 0, Inches(0.12), Inches(1.1), ACCENT)
    add_text(slide, title_text,
             Inches(0.25), Inches(0.12), Inches(11), Inches(0.65),
             font_size=28, bold=True, color=ACCENT)
    if sub_text:
        add_text(slide, sub_text,
                 Inches(0.25), Inches(0.72), Inches(11), Inches(0.35),
                 font_size=14, color=LIGHT_GRY)


def footer(slide, page_num, total=14):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), RGBColor(0x05, 0x28, 0x45))
    add_text(slide, "Log Monitor MLOps  ·  Raúl Brito (A22309632)  ·  IPLUSO 2026",
             Inches(0.2), Inches(7.18), Inches(10), Inches(0.28),
             font_size=9, color=LIGHT_GRY)
    add_text(slide, f"{page_num}/{total}",
             Inches(12.6), Inches(7.18), Inches(0.7), Inches(0.28),
             font_size=9, color=LIGHT_GRY, align=PP_ALIGN.RIGHT)


def content_area():
    """Return (x, y, w, h) for main content below header."""
    return Inches(0.4), Inches(1.25), Inches(12.53), Inches(5.65)


def two_col(slide, left_items, right_items,
            left_title="", right_title="",
            font_size=14, bullet="▸ "):
    cx, cy, cw, ch = content_area()
    col_w = cw / 2 - Inches(0.15)
    col_gap = Inches(0.3)

    for col_idx, (title_str, items) in enumerate(
            [(left_title, left_items), (right_title, right_items)]):
        cx_col = cx + col_idx * (col_w + col_gap)
        # column card
        add_rect(slide, cx_col, cy, col_w, ch, RGBColor(0x0A, 0x25, 0x3A))
        if title_str:
            add_rect(slide, cx_col, cy, col_w, Inches(0.4), ACCENT)
            add_text(slide, title_str,
                     cx_col + Inches(0.1), cy + Inches(0.04),
                     col_w - Inches(0.2), Inches(0.35),
                     font_size=14, bold=True, color=DARK_BG)
            item_y = cy + Inches(0.5)
        else:
            item_y = cy + Inches(0.15)

        for item in items:
            add_text(slide, bullet + item,
                     cx_col + Inches(0.15), item_y,
                     col_w - Inches(0.3), Inches(0.45),
                     font_size=font_size, color=WHITE)
            item_y += Inches(0.45)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
# decorative gradient strip top
add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), ACCENT)
# decorative gradient strip bottom
add_rect(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), ACCENT)

# large title
add_text(slide, "LOG MONITOR MLOps",
         Inches(0.8), Inches(1.6), Inches(11.73), Inches(1.4),
         font_size=54, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, "Sistema Híbrido de Deteção de Anomalias em Logs",
         Inches(0.8), Inches(3.0), Inches(11.73), Inches(0.7),
         font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide,
         "Relatório Intermédio  ·  Projeto Integrado LEIA",
         Inches(0.8), Inches(3.7), Inches(11.73), Inches(0.5),
         font_size=16, color=LIGHT_GRY, align=PP_ALIGN.CENTER)

# divider
add_rect(slide, Inches(3.5), Inches(4.4), Inches(6.33), Inches(0.04), ACCENT2)

add_text(slide, "Raúl Brito   |   A22309632   |   Orientador: Acácio Carmona",
         Inches(0.8), Inches(4.55), Inches(11.73), Inches(0.45),
         font_size=15, color=LIGHT_GRY, align=PP_ALIGN.CENTER)

add_text(slide, "IPLUSO — EET  ·  06 / 05 / 2026",
         Inches(0.8), Inches(5.05), Inches(11.73), Inches(0.4),
         font_size=13, color=LIGHT_GRY, align=PP_ALIGN.CENTER)

# Semana 14 badge
add_rect(slide, Inches(5.2), Inches(5.7), Inches(2.9), Inches(0.55), ACCENT)
add_text(slide, "Semana 14  ·  Goals A–C Concluídos",
         Inches(5.2), Inches(5.72), Inches(2.9), Inches(0.5),
         font_size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

footer(slide, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "Agenda", "15 minutos  ·  14 slides")

items = [
    ("01", "O Problema"),
    ("02", "Solução Proposta"),
    ("03", "Arquitectura (10 microserviços)"),
    ("04", "Pipeline ML Híbrido"),
    ("05", "Resultados de Performance"),
    ("06", "Monitorização & Observabilidade"),
    ("07", "Security Hardening (S14)"),
    ("08", "CI/CD & Qualidade"),
    ("09", "Benchmarking & Gap Analysis"),
    ("10", "Pertinência e Viabilidade"),
    ("11", "Planeamento — Goals A–D"),
    ("12", "Estado Atual & Próximos Passos"),
    ("13", "Conclusão"),
]

cx, cy = Inches(0.5), Inches(1.35)
col_w, row_h = Inches(5.9), Inches(0.38)
for i, (num, label) in enumerate(items):
    col = i % 2
    row = i // 2
    x = cx + col * Inches(6.7)
    y = cy + row * row_h
    add_rect(slide, x, y + Inches(0.04), Inches(0.38), Inches(0.3), ACCENT)
    add_text(slide, num, x + Inches(0.04), y + Inches(0.04),
             Inches(0.32), Inches(0.3), font_size=11, bold=True,
             color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + Inches(0.44), y,
             Inches(5.4), Inches(0.38), font_size=14, color=WHITE)

footer(slide, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — O PROBLEMA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "1 — O Problema", "Contexto e Motivação")

# Left: problem statement
cx, cy, cw, ch = content_area()
half = cw / 2 - Inches(0.15)

add_rect(slide, cx, cy, half, ch, RGBColor(0x0A, 0x25, 0x3A))
add_text(slide, "Contexto",
         cx + Inches(0.15), cy + Inches(0.1),
         half - Inches(0.3), Inches(0.35),
         font_size=15, bold=True, color=ACCENT)

problems = [
    "Organizações geram milhões de logs/dia",
    "Análise manual é inviável à escala atual",
    "Ataques modernos requerem deteção em tempo real",
    "GDPR Art. 22 & AI Act exigem explicabilidade",
    "NIS2: requisitos de cibersegurança obrigatórios",
]
for i, p in enumerate(problems):
    add_text(slide, "▸ " + p,
             cx + Inches(0.2), cy + Inches(0.55 + i * 0.46),
             half - Inches(0.35), Inches(0.44),
             font_size=13, color=WHITE)

# Right: SIEM cost problem
rx = cx + half + Inches(0.3)
add_rect(slide, rx, cy, half, ch, RGBColor(0x0A, 0x25, 0x3A))
add_text(slide, "Lacuna de Mercado",
         rx + Inches(0.15), cy + Inches(0.1),
         half - Inches(0.3), Inches(0.35),
         font_size=15, bold=True, color=ORANGE)

gaps = [
    ("€30k–€100k/ano", "Custo das soluções SIEM enterprise"),
    ("Splunk / QRadar", "Desenhados para grandes orgs"),
    ("Sem explicabilidade", "Caixas-negras sem justificação"),
    ("Sem feedback ML", "Regras estáticas, sem aprendizagem"),
    ("Open-source?", "Não há alternativa academicamente rigorosa"),
]
for i, (bold_part, rest) in enumerate(gaps):
    add_text(slide, bold_part + "  ",
             rx + Inches(0.2), cy + Inches(0.55 + i * 0.46),
             Inches(1.6), Inches(0.44),
             font_size=13, bold=True, color=ORANGE)
    add_text(slide, rest,
             rx + Inches(1.85), cy + Inches(0.55 + i * 0.46),
             half - Inches(2.1), Inches(0.44),
             font_size=13, color=LIGHT_GRY)

footer(slide, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUÇÃO PROPOSTA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "2 — Solução Proposta", "Log Monitor MLOps — alternativa open-source, demo-ready")

cx, cy, cw, ch = content_area()

# central tagline
add_text(slide,
         "Sistema híbrido de deteção de anomalias em logs de aplicações web",
         cx, cy, cw, Inches(0.5),
         font_size=17, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# 3 pillars
pillar_w = cw / 3 - Inches(0.1)
pillar_h = Inches(4.0)
pillar_y = cy + Inches(0.65)
pillars = [
    ("Regras SQL\nDeterminísticas", GREEN,
     ["6 padrões de ataque activos",
      "0 falsos positivos",
      "Latência 7,6–40,8 ms/regra",
      "SQL nativo — auditável"]),
    ("Machine Learning\nNão-Supervisionado", ACCENT,
     ["Isolation Forest (novelty)",
      "F1 = 0,838  |  ROC-AUC 0,95",
      "Treino apenas em normais",
      "Detecta ataques nunca vistos"]),
    ("Explicabilidade\nSHAP + LIME", ORANGE,
     ["Cada decisão tem justificação",
      "Conformidade GDPR Art. 22",
      "AI Act: transparência obrigatória",
      "Ablation study confirma ensemble"]),
]
for i, (title, color, items) in enumerate(pillars):
    px = cx + i * (pillar_w + Inches(0.15))
    add_rect(slide, px, pillar_y, pillar_w, pillar_h, RGBColor(0x0A, 0x25, 0x3A))
    add_rect(slide, px, pillar_y, pillar_w, Inches(0.06), color)
    add_text(slide, title,
             px + Inches(0.1), pillar_y + Inches(0.12),
             pillar_w - Inches(0.2), Inches(0.8),
             font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(slide, "▸ " + item,
                 px + Inches(0.12), pillar_y + Inches(1.0 + j * 0.52),
                 pillar_w - Inches(0.24), Inches(0.5),
                 font_size=12, color=WHITE)

footer(slide, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARQUITECTURA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "3 — Arquitectura", "10 microserviços · 5 camadas · Docker Compose")

cx, cy, cw, ch = content_area()

layers = [
    ("Aplicação",        ACCENT,
     "Flask app (gerador de logs + API ingestão)  ·  Pydantic  ·  Flask-Limiter"),
    ("Armazenamento",    RGBColor(0x48, 0xCA, 0xE4),
     "PostgreSQL + TimescaleDB (hypertables)  ·  ACID  ·  séries temporais"),
    ("Detecção",         GREEN,
     "Rule Engine SQL (6 regras)  +  ML Pipeline (Isolation Forest + Random Forest + MLflow)"),
    ("Explicabilidade",  ORANGE,
     "SHAP  ·  LIME  ·  Ablation study  ·  Feedback loop"),
    ("Monitorização",    RGBColor(0xFF, 0xD1, 0x66),
     "Prometheus (9 métricas custom)  ·  Grafana (3 dashboards)  ·  Streamlit dashboard  ·  SLOs"),
    ("Infraestrutura",   LIGHT_GRY,
     "Docker Compose (multi-stage builds: ~1,2 GB → ~180 MB/img)  ·  GitHub Actions CI/CD"),
]

layer_h = Inches(0.68)
for i, (name, color, desc) in enumerate(layers):
    y = cy + i * (layer_h + Inches(0.06))
    add_rect(slide, cx, y, Inches(2.0), layer_h, color)
    add_text(slide, name,
             cx + Inches(0.05), y + Inches(0.15),
             Inches(1.9), Inches(0.4),
             font_size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, cx + Inches(2.05), y, cw - Inches(2.05), layer_h,
             RGBColor(0x0A, 0x25, 0x3A))
    add_text(slide, desc,
             cx + Inches(2.2), y + Inches(0.16),
             cw - Inches(2.35), Inches(0.4),
             font_size=12.5, color=WHITE)

footer(slide, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PIPELINE ML HÍBRIDO
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "4 — Pipeline ML Híbrido", "Isolation Forest + Random Forest + SHAP/LIME")

cx, cy, cw, ch = content_area()
col_w = cw / 2 - Inches(0.15)

# Left — Isolation Forest
add_rect(slide, cx, cy, col_w, ch, RGBColor(0x0A, 0x25, 0x3A))
add_rect(slide, cx, cy, col_w, Inches(0.38), ACCENT)
add_text(slide, "Isolation Forest (não-supervisionado)",
         cx + Inches(0.1), cy + Inches(0.04), col_w - Inches(0.2), Inches(0.35),
         font_size=13, bold=True, color=DARK_BG)

if_items = [
    "Treino: apenas logs normais",
    "Teste: ataques nunca vistos",
    "F1-Score = 0,838",
    "Precision@1% = 0,941",
    "ROC-AUC = 0,950",
    "Protocolo novelty-by-scenario",
    "→ detecta zero-day sem labels",
]
for i, item in enumerate(if_items):
    clr = GREEN if item.startswith("F1") or item.startswith("Pr") or item.startswith("ROC") else WHITE
    bold = clr == GREEN
    add_text(slide, ("▸ " if not item.startswith("→") else "") + item,
             cx + Inches(0.15), cy + Inches(0.5 + i * 0.52),
             col_w - Inches(0.3), Inches(0.5),
             font_size=13, bold=bold, color=clr)

# Right — Random Forest supervised + ensemble
rx = cx + col_w + Inches(0.3)
add_rect(slide, rx, cy, col_w, ch, RGBColor(0x0A, 0x25, 0x3A))
add_rect(slide, rx, cy, col_w, Inches(0.38), ORANGE)
add_text(slide, "Random Forest (supervisionado) + Ensemble",
         rx + Inches(0.1), cy + Inches(0.04), col_w - Inches(0.2), Inches(0.35),
         font_size=13, bold=True, color=DARK_BG)

rf_items = [
    "Treino com labels de ataque",
    "F1-Score = 0,783",
    "Ablation: ensemble > IF solo",
    "SHAP: importância de features",
    "LIME: explicações locais",
    "MLflow: tracking & versionamento",
    "→ conformidade GDPR Art. 22",
]
for i, item in enumerate(rf_items):
    clr = ORANGE if item.startswith("F1") or item.startswith("Ab") else WHITE
    bold = clr == ORANGE
    add_text(slide, ("▸ " if not item.startswith("→") else "") + item,
             rx + Inches(0.15), cy + Inches(0.5 + i * 0.52),
             col_w - Inches(0.3), Inches(0.5),
             font_size=13, bold=bold, color=clr)

footer(slide, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RESULTADOS DE PERFORMANCE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "5 — Resultados de Performance", "Métricas medidas em testes reais")

cx, cy, cw, ch = content_area()

# 4 metric cards
card_w = cw / 2 - Inches(0.12)
card_h = ch / 2 - Inches(0.1)
metrics = [
    ("Ingestão de Logs",   ACCENT,
     "10.933 logs/s  (batch=500)",
     ["Batch=50:  6.493 logs/s",
      "Batch=100: 7.145 logs/s  (+10%)",
      "Batch=500: 10.933 logs/s (+68%)",
      "Batch=1000: 13.646 logs/s (+110%)",
      "100.000 logs em 9 s  ·  39M/hora"]),
    ("Testes de Carga",    GREEN,
     "50 users  ·  0 falhas",
     ["Latência mediana: 3 ms",
      "p95: 10 ms  ·  p99: 25 ms",
      "50.000 logs stress test: OK",
      "Locust 30s run, 50 users",
      "0% error rate"]),
    ("Rule Engine SQL",    ORANGE,
     "6 regras  ·  0 falsos positivos",
     ["Latência por regra: 7,6–40,8 ms",
      "Brute force, scanning, SQLi...",
      "Rate abuse, DDoS, path traversal",
      "Deteção em tempo real",
      "0 false positives confirmados"]),
    ("Qualidade de Código", RGBColor(0xFF, 0xD1, 0x66),
     "Pylint 9,18/10  ·  81 testes",
     ["Cobertura: 70,97% (CI ≥70% ✓)",
      "0 HIGH/MEDIUM no Bandit",
      "54 testes de segurança OK",
      "CI threshold cumprido",
      "5 containers non-root"]),
]
for i, (title, color, headline, bullets) in enumerate(metrics):
    col = i % 2
    row = i // 2
    mx = cx + col * (card_w + Inches(0.24))
    my = cy + row * (card_h + Inches(0.2))
    add_rect(slide, mx, my, card_w, card_h, RGBColor(0x0A, 0x25, 0x3A))
    add_rect(slide, mx, my, card_w, Inches(0.36), color)
    add_text(slide, title,
             mx + Inches(0.1), my + Inches(0.03), card_w - Inches(0.2), Inches(0.32),
             font_size=13, bold=True, color=DARK_BG)
    add_text(slide, headline,
             mx + Inches(0.1), my + Inches(0.40), card_w - Inches(0.2), Inches(0.35),
             font_size=13, bold=True, color=color)
    for j, b in enumerate(bullets):
        add_text(slide, "▸ " + b,
                 mx + Inches(0.12), my + Inches(0.78 + j * 0.38),
                 card_w - Inches(0.24), Inches(0.36),
                 font_size=11, color=LIGHT_GRY)

footer(slide, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MONITORIZAÇÃO
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "6 — Monitorização & Observabilidade",
           "Prometheus + Grafana + Streamlit + SLOs definidos")

cx, cy, cw, ch = content_area()

# 3 columns
col_w = cw / 3 - Inches(0.1)
cols = [
    ("Prometheus", ACCENT,
     ["9 métricas custom logmonitor_*",
      "Scrape interval: 15 s",
      "Alertmanager configurado",
      "Integração nativa com Flask",
      "Retenção 15 dias"]),
    ("Grafana", GREEN,
     ["3 dashboards auto-provisionados",
      "Overview · ML · Security",
      "Provisioning via YAML",
      "SLO visibility em tempo real",
      "Zero-config ao startup"]),
    ("SLOs Definidos", ORANGE,
     ["Disponibilidade ≥ 99,5%",
      "p95 latência < 200 ms",
      "F1 do modelo ≥ 0,75",
      "Data freshness < 5 min",
      "Streamlit dashboard operacional"]),
]
for i, (title, color, items) in enumerate(cols):
    x = cx + i * (col_w + Inches(0.15))
    add_rect(slide, x, cy, col_w, ch, RGBColor(0x0A, 0x25, 0x3A))
    add_rect(slide, x, cy, col_w, Inches(0.4), color)
    add_text(slide, title,
             x + Inches(0.08), cy + Inches(0.04), col_w - Inches(0.16), Inches(0.34),
             font_size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(slide, "▸ " + item,
                 x + Inches(0.12), cy + Inches(0.52 + j * 0.54),
                 col_w - Inches(0.24), Inches(0.5),
                 font_size=12.5, color=WHITE)

footer(slide, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SECURITY HARDENING
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "7 — Security Hardening (S14)", "Segurança aplicada em toda a stack")

cx, cy, cw, ch = content_area()
col_w = cw / 2 - Inches(0.15)

# Left: measures
add_rect(slide, cx, cy, col_w, ch, RGBColor(0x0A, 0x25, 0x3A))
add_text(slide, "Medidas Implementadas",
         cx + Inches(0.1), cy + Inches(0.08), col_w - Inches(0.2), Inches(0.35),
         font_size=15, bold=True, color=GREEN)
sec_items = [
    "5 Dockerfiles: non-root user (appuser)",
    "Rate limiting em rotas sensíveis",
    "Validação de inputs (Pydantic)",
    "Bandit: 0 HIGH / 0 MEDIUM findings",
    "Trivy container scanning no CI",
    "Dependabot: dependências actualizadas",
    "54 testes de segurança a passar",
    "Secrets em variáveis de ambiente",
]
for i, item in enumerate(sec_items):
    add_text(slide, "✓ " + item,
             cx + Inches(0.15), cy + Inches(0.55 + i * 0.48),
             col_w - Inches(0.3), Inches(0.44),
             font_size=12.5, color=WHITE)

# Right: context
rx = cx + col_w + Inches(0.3)
add_rect(slide, rx, cy, col_w, ch, RGBColor(0x0A, 0x25, 0x3A))
add_text(slide, "Contexto Regulatório",
         rx + Inches(0.1), cy + Inches(0.08), col_w - Inches(0.2), Inches(0.35),
         font_size=15, bold=True, color=ORANGE)

regs = [
    ("OWASP Top 10", "Input validation, rate limiting, secrets"),
    ("GDPR Art. 22", "Explicabilidade de decisões automáticas"),
    ("AI Act (2024)", "Sistemas de IA de alto risco: transparência"),
    ("NIS2", "Cibersegurança obrigatória para serviços essenciais"),
    ("Non-root containers", "Princípio do menor privilégio"),
]
for i, (reg, desc) in enumerate(regs):
    y = cy + Inches(0.55 + i * 0.72)
    add_rect(slide, rx + Inches(0.1), y, col_w - Inches(0.2), Inches(0.6),
             RGBColor(0x12, 0x30, 0x4A))
    add_text(slide, reg,
             rx + Inches(0.2), y + Inches(0.04), col_w - Inches(0.4), Inches(0.28),
             font_size=12, bold=True, color=ORANGE)
    add_text(slide, desc,
             rx + Inches(0.2), y + Inches(0.3), col_w - Inches(0.4), Inches(0.28),
             font_size=11, color=LIGHT_GRY)

footer(slide, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CI/CD & QUALIDADE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "8 — CI/CD & Qualidade", "GitHub Actions · 81 testes · Cobertura 70,97%")

cx, cy, cw, ch = content_area()

# Pipeline flow
stages = [
    ("Push\nGit", ACCENT),
    ("Pylint\n9,18/10", GREEN),
    ("Bandit\n0 HIGH", GREEN),
    ("Trivy\nScanning", GREEN),
    ("Unit Tests\n70,97%", GREEN),
    ("Load Test\n0 falhas", GREEN),
    ("Deploy\nLocal", ACCENT),
]
stage_w = cw / len(stages) - Inches(0.05)
stage_h = Inches(1.1)
sy = cy + Inches(0.3)
for i, (label, color) in enumerate(stages):
    sx = cx + i * (stage_w + Inches(0.05))
    add_rect(slide, sx, sy, stage_w, stage_h, color)
    add_text(slide, label,
             sx + Inches(0.03), sy + Inches(0.1),
             stage_w - Inches(0.06), Inches(0.9),
             font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        add_text(slide, "→",
                 sx + stage_w, sy + Inches(0.35),
                 Inches(0.05), Inches(0.4),
                 font_size=14, bold=True, color=WHITE)

# Details below
detail_y = cy + Inches(1.65)
detail_items = [
    ("81 testes\nautomatizados", ACCENT),
    ("Cobertura\n70,97% ✓", GREEN),
    ("Pylint\n9,18/10", GREEN),
    ("Threshold CI\n≥ 70%", ORANGE),
    ("Dependabot\nautomático", ACCENT2),
    ("Multi-stage\nbuild Docker", LIGHT_GRY),
    ("Semana 13\nIntegrado", WHITE),
]
dw = cw / len(detail_items) - Inches(0.05)
dh = Inches(1.1)
for i, (label, color) in enumerate(detail_items):
    dx = cx + i * (dw + Inches(0.05))
    add_rect(slide, dx, detail_y, dw, dh, RGBColor(0x0A, 0x25, 0x3A))
    add_rect(slide, dx, detail_y, dw, Inches(0.06), color)
    add_text(slide, label,
             dx + Inches(0.05), detail_y + Inches(0.12),
             dw - Inches(0.1), Inches(0.9),
             font_size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

# Bottom note
add_text(slide,
         "Pipeline CI/CD garante qualidade e segurança em cada commit — nenhum merge sem testes passados",
         cx, detail_y + Inches(1.2), cw, Inches(0.45),
         font_size=13, color=LIGHT_GRY, italic=True, align=PP_ALIGN.CENTER)

footer(slide, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — BENCHMARKING
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "9 — Benchmarking & Gap Analysis",
           "Posicionamento face a soluções comerciais")

cx, cy, cw, ch = content_area()

# Table-style comparison
headers = ["Solução", "Custo/ano", "Explicabilidade", "ML Adaptativo", "Open-source"]
rows = [
    ("Splunk SIEM",     "€50k–€100k", "Não", "Parcial", "Não"),
    ("IBM QRadar",      "€30k–€80k",  "Não", "Limitado", "Não"),
    ("Elastic SIEM",    "€10k–€30k",  "Parcial", "Não", "Parcial"),
    ("Datadog",         "€20k–€60k",  "Não", "Não", "Não"),
    ("Log Monitor MLOps","€0 (OS)", "SHAP+LIME", "Sim", "Sim"),
]
col_widths = [Inches(2.5), Inches(1.7), Inches(2.0), Inches(2.1), Inches(1.8)]
col_x = [cx]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w + Inches(0.05))

row_h = Inches(0.52)
# header
for j, (hdr, w) in enumerate(zip(headers, col_widths)):
    add_rect(slide, col_x[j], cy, w, Inches(0.38), ACCENT)
    add_text(slide, hdr,
             col_x[j] + Inches(0.05), cy + Inches(0.05),
             w - Inches(0.1), Inches(0.3),
             font_size=12, bold=True, color=DARK_BG)

for i, row in enumerate(rows):
    is_our = "Log Monitor" in row[0]
    row_color = RGBColor(0x05, 0x38, 0x2A) if is_our else RGBColor(0x0A, 0x25, 0x3A)
    ry = cy + Inches(0.45) + i * row_h
    for j, (cell, w) in enumerate(zip(row, col_widths)):
        add_rect(slide, col_x[j], ry, w, row_h - Inches(0.04), row_color)
        clr = GREEN if is_our else (
            RGBColor(0xFF, 0x5C, 0x5C) if cell in ("Não", "Limitado", "Parcial") else WHITE)
        add_text(slide, cell,
                 col_x[j] + Inches(0.07), ry + Inches(0.1),
                 w - Inches(0.14), row_h - Inches(0.14),
                 font_size=12, bold=is_our, color=clr)

add_text(slide,
         "Gap preenchido: explicabilidade nativa, ML adaptativo, open-source — a €0 vs €30k–€100k/ano",
         cx, cy + Inches(3.2), cw, Inches(0.45),
         font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

footer(slide, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PERTINÊNCIA E VIABILIDADE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "10 — Pertinência e Viabilidade",
           "Regulatória · Científica · Económica")

cx, cy, cw, ch = content_area()
col_w = cw / 3 - Inches(0.12)

dimensions = [
    ("Regulatória", ACCENT,
     ["GDPR Art. 22: direito à explicação",
      "AI Act (2024): sistemas de alto risco",
      "NIS2: cibersegurança obrigatória",
      "Conformidade por design"]),
    ("Científica", GREEN,
     ["XAI aplicado a cibersegurança",
      "Ablation study: ensemble > solo",
      "Novelty detection sem labels",
      "Publicável — contribuição real"]),
    ("Económica", ORANGE,
     ["254 h de trabalho individual",
      "€0 custo de ferramentas",
      "Equivale a €30k–€100k/ano SIEM",
      "Viável para PMEs e academia"]),
]
for i, (title, color, items) in enumerate(dimensions):
    x = cx + i * (col_w + Inches(0.18))
    add_rect(slide, x, cy, col_w, ch, RGBColor(0x0A, 0x25, 0x3A))
    add_rect(slide, x, cy, col_w, Inches(0.06), color)
    # icon-like number
    add_rect(slide, x + col_w/2 - Inches(0.3), cy + Inches(0.1),
             Inches(0.6), Inches(0.6), color)
    add_text(slide, str(i+1),
             x + col_w/2 - Inches(0.3), cy + Inches(0.14),
             Inches(0.6), Inches(0.5),
             font_size=20, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(slide, title,
             x + Inches(0.1), cy + Inches(0.8),
             col_w - Inches(0.2), Inches(0.4),
             font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(slide, "▸ " + item,
                 x + Inches(0.12), cy + Inches(1.3 + j * 0.62),
                 col_w - Inches(0.24), Inches(0.58),
                 font_size=12.5, color=WHITE)

footer(slide, 12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — PLANEAMENTO GOALS A–D
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "11 — Planeamento — Goals A–D",
           "27 semanas  ·  4 fases incrementais  ·  4 milestones")

cx, cy, cw, ch = content_area()

goals = [
    ("Goal A", "S1–S8", ACCENT,
     "Setup + Ingestão + Rule Engine + ML + Explainability", "Concluído ✓"),
    ("Goal B", "S9–S13", GREEN,
     "Docker Orchestration + Monitorização + Dashboard + Testing + CI/CD", "Concluído ✓"),
    ("Goal C", "S14–S16", ORANGE,
     "Security Hardening + Demo + Documentação Intercalar", "Em curso"),
    ("Goal D", "S17–S23", RGBColor(0xFF, 0xD1, 0x66),
     "API REST Pública + Alertas + Dashboard Avançado + Relatório Final", "Planeado"),
]

goal_h = ch / len(goals) - Inches(0.08)
for i, (goal, weeks, color, desc, status) in enumerate(goals):
    gy = cy + i * (goal_h + Inches(0.08))
    # left badge
    add_rect(slide, cx, gy, Inches(1.2), goal_h, color)
    add_text(slide, goal,
             cx, gy + Inches(0.08), Inches(1.2), Inches(0.38),
             font_size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(slide, weeks,
             cx, gy + Inches(0.48), Inches(1.2), Inches(0.3),
             font_size=12, color=DARK_BG, align=PP_ALIGN.CENTER)
    # description
    add_rect(slide, cx + Inches(1.25), gy, cw - Inches(2.65), goal_h,
             RGBColor(0x0A, 0x25, 0x3A))
    add_text(slide, desc,
             cx + Inches(1.4), gy + Inches(0.2),
             cw - Inches(3.0), goal_h - Inches(0.4),
             font_size=13, color=WHITE)
    # status badge
    s_color = GREEN if "Concluído" in status else (ORANGE if "curso" in status else LIGHT_GRY)
    add_rect(slide, cx + cw - Inches(1.35), gy + Inches(0.12),
             Inches(1.3), Inches(0.45), s_color)
    add_text(slide, status,
             cx + cw - Inches(1.35), gy + Inches(0.15),
             Inches(1.3), Inches(0.38),
             font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

footer(slide, 13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — ESTADO ACTUAL & PRÓXIMOS PASSOS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
header_bar(slide, "12 — Estado Actual & Próximos Passos", "Semana 14 — 100% operacional")

cx, cy, cw, ch = content_area()
col_w = cw / 2 - Inches(0.15)

# Left: current state
add_rect(slide, cx, cy, col_w, ch, RGBColor(0x0A, 0x25, 0x3A))
add_text(slide, "Estado Actual (S14) — Operacional",
         cx + Inches(0.1), cy + Inches(0.08), col_w - Inches(0.2), Inches(0.38),
         font_size=14, bold=True, color=GREEN)
current = [
    "Ingestão: 10.933 logs/s (batch=500)",
    "Rule Engine: 6 regras, 0 FP",
    "IF: F1=0,838  |  RF: F1=0,783",
    "81 testes, cobertura 70,97% ✓",
    "3 dashboards Grafana ativos",
    "CI/CD: Pylint, Bandit, Trivy OK",
    "Security: 0 HIGH/MEDIUM findings",
    "5 containers non-root",
]
for i, item in enumerate(current):
    add_text(slide, "✓ " + item,
             cx + Inches(0.15), cy + Inches(0.56 + i * 0.47),
             col_w - Inches(0.3), Inches(0.44),
             font_size=12.5, color=WHITE)

# Right: next steps
rx = cx + col_w + Inches(0.3)
add_rect(slide, rx, cy, col_w, ch, RGBColor(0x0A, 0x25, 0x3A))
add_text(slide, "Próximas Semanas (S15–S16)",
         rx + Inches(0.1), cy + Inches(0.08), col_w - Inches(0.2), Inches(0.38),
         font_size=14, bold=True, color=ACCENT)
next_items = [
    ("S15", "Cobertura testes ≥ 75%"),
    ("S15", "Documentação intercalar"),
    ("S16", "Demo ao vivo preparada"),
    ("S16", "Relatório final revisto"),
    ("S17+", "Goal D: API REST pública"),
    ("S17+", "Alerting avançado"),
    ("S23", "Relatório final completo"),
]
for i, (week, item) in enumerate(next_items):
    y = cy + Inches(0.56 + i * 0.47)
    add_rect(slide, rx + Inches(0.1), y, Inches(0.48), Inches(0.38),
             ACCENT if "S15" in week or "S16" in week else RGBColor(0x30, 0x50, 0x70))
    add_text(slide, week,
             rx + Inches(0.1), y + Inches(0.05), Inches(0.48), Inches(0.3),
             font_size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(slide, item,
             rx + Inches(0.65), y, col_w - Inches(0.8), Inches(0.44),
             font_size=12.5, color=WHITE)

footer(slide, 14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CONCLUSÃO
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), ACCENT)
add_rect(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), ACCENT)

add_text(slide, "Conclusão",
         Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.7),
         font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(1.0), Inches(1.3), Inches(11.33), Inches(0.04), ACCENT2)

conclusions = [
    "Sistema 100% operacional em ambiente containerizado — demonstração ao vivo possível",
    "Goals A, B e C concluídos em 14 semanas com desenvolvimento individual",
    "Performance comprovada: 10.933 logs/s, F1=0,838, 0 falhas em stress test",
    "Segurança robusta: non-root containers, 0 HIGH/MEDIUM Bandit, 54 testes de segurança",
    "Explicabilidade nativa (SHAP + LIME) em conformidade com GDPR Art. 22 e AI Act",
    "Alternativa open-source a soluções €30k–€100k/ano — viável para PMEs e academia",
]
for i, c in enumerate(conclusions):
    add_text(slide, "▸  " + c,
             Inches(1.0), Inches(1.55 + i * 0.62),
             Inches(11.33), Inches(0.58),
             font_size=15, color=WHITE)

add_rect(slide, Inches(1.5), Inches(5.5), Inches(10.33), Inches(0.04), ACCENT2)

add_text(slide, "Obrigado  ·  Questões?",
         Inches(0.8), Inches(5.65), Inches(11.73), Inches(0.55),
         font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, "Raúl Brito (A22309632)  ·  raulrazer123@gmail.com  ·  IPLUSO — EET  ·  2026",
         Inches(0.8), Inches(6.25), Inches(11.73), Inches(0.4),
         font_size=13, color=LIGHT_GRY, align=PP_ALIGN.CENTER)

footer(slide, 15, 15)


# ── Save ───────────────────────────────────────────────────────────────────────
out_path = "docs/Relatorio/LogMonitorMLOps_Apresentacao_v1.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
