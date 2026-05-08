"""Enhance LogMonitorMLOps_Apresentacao_v2.pptx with real screenshots.

Strategy:
- Slide 5 (Architecture): wide DockerCompose PS banner at bottom (OK, wide image)
- 3 new dedicated evidence slides (large images, 2 per slide or 1 large):
    NEW-A after ML Pipeline (index 5): mlFlowMEtric — single large image
    NEW-B after Monitoring (index 8): Grafana.png + Prometheus.png side by side
    NEW-C after NEW-B:                StreamLit.png + AlertManager.png side by side
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
import os

SRC = "docs/Relatorio/LogMonitorMLOps_Apresentacao_v2.pptx"
DST = "docs/Relatorio/LogMonitorMLOps_Apresentacao_v3.pptx"
IMG = "docs/report_media/ParaApresentacao/"

prs = Presentation(SRC)


def add_img(slide, fname, left, top, width=None, height=None):
    path = os.path.join(IMG, fname)
    return slide.shapes.add_picture(path, left, top, width=width, height=height)


def clone_slide(prs, source_idx):
    """Deep-copy a slide and append it at the end of the presentation."""
    src = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    new_spTree = new_slide.shapes._spTree
    for child in list(new_spTree):
        new_spTree.remove(child)
    for child in src.shapes._spTree:
        new_spTree.append(copy.deepcopy(child))
    return new_slide


def clear_content_shapes(slide, keep_names):
    """Remove all shapes whose name is not in keep_names."""
    to_remove = [s._element for s in slide.shapes if s.name not in keep_names]
    for el in to_remove:
        el.getparent().remove(el)


def set_first_run(slide, shape_name, new_text):
    """Replace text in the first non-empty run of every paragraph in a shape."""
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.runs:
                    para.runs[0].text = new_text
                    for r in para.runs[1:]:
                        r.text = ""
            return


def update_slide_number(slide, num, total):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if "/" in t and len(t) <= 6 and t[0].isdigit():
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "/" in run.text:
                        run.text = f"{num}/{total}"
                        return


# ── Slide 5 (index 4): Architecture ── DockerCompose PS banner ───────────────
# Tall narrow terminal shot (ratio 7:1) — perfect as a wide thin banner
slide5 = prs.slides[4]
add_img(slide5, "2DockerComposePS.png",
        left=Inches(0.40), top=Inches(5.80),
        width=Inches(12.53))  # height auto ≈ 1.79"

# ── Create 3 new evidence slides (appended at end, then repositioned) ─────────

# Background/header/footer shape names to keep when clearing a cloned slide
KEEP = {"Rectangle 1", "Rectangle 2", "Rectangle 3",
        "TextBox 4", "TextBox 5",
        "Rectangle 30", "TextBox 31", "TextBox 32"}

# Source for cloning: slide 9 (Monitoring, index 8) — simple blank+header design
CLONE_SRC = 8

# ── NEW-A: MLflow Evidence (goes after ML Pipeline, index 5) ─────────────────
slide_a = clone_slide(prs, CLONE_SRC)
clear_content_shapes(slide_a, KEEP)
set_first_run(slide_a, "TextBox 4", "MLflow — Experiment Tracking")
set_first_run(slide_a, "TextBox 5",
              "Rastreamento de experimentos · Versionamento de modelos · Métricas ao vivo")
# Single large image, centered in content area
# Content area: top=1.15" to 7.00" = 5.85" h, full width ~13.33"
# mlFlowMEtric ratio ≈ 1.98 → at width=11.0": height=5.56"
add_img(slide_a, "mlFlowMEtric.png",
        left=Inches(1.17), top=Inches(1.35),
        width=Inches(11.0))  # height auto ≈ 5.56"

# ── NEW-B: Grafana + Prometheus (goes after Monitoring, index 8) ──────────────
slide_b = clone_slide(prs, CLONE_SRC)
clear_content_shapes(slide_b, KEEP)
set_first_run(slide_b, "TextBox 4", "Grafana + Prometheus — Dashboards ao Vivo")
set_first_run(slide_b, "TextBox 5",
              "3 dashboards auto-provisionados · 9 métricas custom · scrape 15 s")
# Two large images side by side
# Each image ≈ 6.30" wide; at ratio 1.97 → height ≈ 3.20"
# Vertically centered: content area 5.85", image 3.20" → top offset = (5.85-3.20)/2 = 1.33"
# top = 1.15 + 1.33 = 2.48 → round to 2.15 for breathing room near title
IMG_TOP = Inches(2.10)
IMG_W = Inches(6.30)
add_img(slide_b, "Grafana.png",   left=Inches(0.22), top=IMG_TOP, width=IMG_W)
add_img(slide_b, "Prometheus.png", left=Inches(6.81), top=IMG_TOP, width=IMG_W)

# ── NEW-C: Streamlit + AlertManager (goes right after NEW-B) ─────────────────
slide_c = clone_slide(prs, CLONE_SRC)
clear_content_shapes(slide_c, KEEP)
set_first_run(slide_c, "TextBox 4", "Streamlit + Alertmanager — Monitorização")
set_first_run(slide_c, "TextBox 5",
              "Dashboard interativo · Alertas configurados · SLOs visíveis em tempo real")
add_img(slide_c, "StreamLit.png",    left=Inches(0.22), top=IMG_TOP, width=IMG_W)
add_img(slide_c, "AlertManager.png", left=Inches(6.81), top=IMG_TOP, width=IMG_W)

# ── Reorder slides ─────────────────────────────────────────────────────────────
# After the three clones the deck is:
#   indices 0-15: original slides
#   index 16: slide_a (MLflow)
#   index 17: slide_b (Grafana+Prom)
#   index 18: slide_c (Streamlit+Alert)
#
# Target order:
#   0-5, NEW-A, 6-8, NEW-B, NEW-C, 9-15
#   i.e. NEW-A after ML Pipeline (5), NEW-B and NEW-C after Monitoring (8)

sld_id_lst = prs.slides._sldIdLst
all_ids = list(sld_id_lst)
id_a = all_ids[16]
id_b = all_ids[17]
id_c = all_ids[18]

# Remove new slides from the end
for elem in (id_a, id_b, id_c):
    sld_id_lst.remove(elem)
# sld_id_lst now has 0-15 (original)

# Insert in reverse order of target position so earlier insertions don't shift later targets
# NEW-C and NEW-B after monitoring (index 8 in original list = position 8)
# Insert NEW-C at 9 first, then NEW-B at 9 (NEW-C shifts to 10)
sld_id_lst.insert(9, id_c)   # after monitoring: ...8, NEW-C, 9...
sld_id_lst.insert(9, id_b)   # before NEW-C:     ...8, NEW-B, NEW-C, 9...
# NEW-A after ML Pipeline (index 5 in original list = position 5)
sld_id_lst.insert(6, id_a)   # after ML Pipeline: ...5, NEW-A, 6...

# Final order:
# 0-5, NEW-A(6), 6→7, 7→8, 8→9, NEW-B(10), NEW-C(11), 9→12, 10→13, ..., 15→18

# ── Update slide number footers ────────────────────────────────────────────────
# 19 total slides; content numbered 1–18 (last two share 18/18 like original)
TOTAL = 18
for i, slide in enumerate(prs.slides):
    num = min(i + 1, TOTAL)
    update_slide_number(slide, num, TOTAL)

# Set the numbers on the 3 new slides explicitly (they were cloned from slide 9 = "9/15")
for slide, expected_idx in [(slide_a, 6), (slide_b, 10), (slide_c, 11)]:
    update_slide_number(slide, min(expected_idx + 1, TOTAL), TOTAL)

prs.save(DST)
print(f"Saved: {DST}  ({len(prs.slides)} slides)")
